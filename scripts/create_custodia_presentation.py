from __future__ import annotations

import io
from pathlib import Path

from PIL import Image, ImageDraw, ImageEnhance, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
BACKGROUND = ROOT / "src" / "assets" / "forensics-background.png"
OUTPUT = ROOT / "Custodia_Group-I_Presentation_v8.pptx"
PRESENTATION_FONT = "Comic Sans MS"

W = Inches(13.333)
H = Inches(7.5)

NAVY = RGBColor(9, 26, 48)
PANEL = RGBColor(16, 47, 79)
PANEL_2 = RGBColor(21, 62, 101)
CYAN = RGBColor(119, 222, 255)
TEAL = RGBColor(116, 235, 198)
WHITE = RGBColor(248, 252, 255)
MUTED = RGBColor(211, 229, 244)
AMBER = RGBColor(246, 184, 89)
RED = RGBColor(247, 126, 137)


def rgb(hex_value: str) -> tuple[int, int, int]:
    hex_value = hex_value.lstrip("#")
    return tuple(int(hex_value[index : index + 2], 16) for index in (0, 2, 4))


def image_bytes(image: Image.Image, fmt: str = "PNG") -> io.BytesIO:
    stream = io.BytesIO()
    image.save(stream, format=fmt)
    stream.seek(0)
    return stream


def make_background() -> io.BytesIO:
    source = Image.open(BACKGROUND).convert("RGB")
    source = ImageOps.fit(source, (1920, 1080), method=Image.Resampling.LANCZOS, centering=(0.5, 0.5))
    source = ImageEnhance.Contrast(source).enhance(1.25)
    source = ImageEnhance.Brightness(source).enhance(0.82)
    tint = Image.new("RGB", source.size, rgb("#081c35"))
    source = Image.blend(source, tint, 0.30)
    return image_bytes(source)


def avatar_bytes(index: int, feminine: bool) -> io.BytesIO:
    size = 520
    canvas = Image.new("RGBA", (size, size), (0, 0, 0, 0))
    draw = ImageDraw.Draw(canvas)
    palettes = [
        ("#2b71d9", "#d7e8ff", "#163b69"),
        ("#2c9d9a", "#f1d0b1", "#3d2133"),
        ("#8758cf", "#a96f4f", "#172338"),
        ("#d17a42", "#f0c39e", "#35213e"),
        ("#2386ba", "#8d5a3d", "#20192c"),
        ("#4c9d64", "#e7b68e", "#40241d"),
    ]
    accent, skin, hair = palettes[index % len(palettes)]
    draw.ellipse((12, 12, size - 12, size - 12), fill=rgb(accent), outline=rgb("#a7e7ff"), width=7)
    draw.ellipse((113, 104, 407, 398), fill=rgb(skin))
    if feminine:
        draw.ellipse((82, 66, 438, 285), fill=rgb(hair))
        draw.ellipse((106, 151, 414, 414), fill=rgb(skin))
        draw.pieslice((80, 240, 440, 540), 0, 180, fill=rgb(hair))
    else:
        draw.pieslice((103, 55, 417, 255), 180, 360, fill=rgb(hair))
        draw.rectangle((114, 145, 406, 250), fill=rgb(hair))
    draw.ellipse((176, 218, 203, 245), fill=rgb("#15263a"))
    draw.ellipse((317, 218, 344, 245), fill=rgb("#15263a"))
    draw.arc((211, 252, 309, 326), 15, 165, fill=rgb("#813f4d"), width=9)
    draw.rectangle((238, 355, 282, 412), fill=rgb(skin))
    draw.ellipse((86, 374, 434, 650), fill=rgb("#0d294b"))
    draw.arc((144, 392, 376, 635), 180, 360, fill=rgb("#6bd7d1"), width=8)
    return image_bytes(canvas)


def make_tools_visual() -> io.BytesIO:
    """Create a small forensic-tool illustration for the presentation."""
    canvas = Image.new("RGB", (900, 520), rgb("#081c35"))
    draw = ImageDraw.Draw(canvas)
    for x in range(0, 900, 90):
        draw.line((x, 0, x - 180, 520), fill=rgb("#12385b"), width=2)
    draw.rounded_rectangle((28, 28, 872, 492), radius=28, outline=rgb("#2d648b"), width=3)
    draw.ellipse((334, 150, 566, 382), fill=rgb("#103c5e"), outline=rgb("#63cdff"), width=5)
    draw.ellipse((392, 205, 508, 321), outline=rgb("#53dcb5"), width=10)
    draw.line((485, 300, 548, 363), fill=rgb("#53dcb5"), width=10)
    draw.arc((407, 220, 493, 302), 205, 335, fill=rgb("#f6b859"), width=5)
    items = [
        (142, 112, "FILE", "#63cdff"),
        (642, 112, "HASH", "#53dcb5"),
        (142, 354, "NODES", "#f6b859"),
        (642, 354, "AUDIT", "#f77e89"),
    ]
    for x, y, label, color in items:
        draw.rounded_rectangle((x - 68, y - 38, x + 68, y + 38), radius=18, fill=rgb("#102f4f"), outline=rgb(color), width=4)
        draw.ellipse((x - 43, y - 18, x - 13, y + 12), outline=rgb(color), width=4)
        draw.text((x - 2, y - 15), label, fill=rgb("#f2f8ff"))
    for start, end in [((210, 112), (350, 176)), ((690, 112), (550, 176)), ((210, 354), (350, 354)), ((690, 354), (550, 354))]:
        draw.line((*start, *end), fill=rgb("#3d789c"), width=4)
    draw.text((302, 438), "forensic evidence toolkit", fill=rgb("#b8cfe5"))
    return image_bytes(canvas)


def make_nodes_visual() -> io.BytesIO:
    """Create a node/quorum illustration for the replication slide."""
    canvas = Image.new("RGB", (900, 520), rgb("#081c35"))
    draw = ImageDraw.Draw(canvas)
    center = (450, 260)
    draw.ellipse((328, 138, 572, 382), fill=rgb("#103c5e"), outline=rgb("#53dcb5"), width=6)
    draw.ellipse((367, 177, 533, 343), outline=rgb("#63cdff"), width=4)
    draw.text((401, 238), "QUORUM", fill=rgb("#f2f8ff"))
    nodes = [(150, 115, "ATLAS", "#63cdff"), (750, 115, "BOREAL", "#53dcb5"), (150, 405, "CINDER", "#f6b859"), (750, 405, "DELTA", "#f77e89")]
    for x, y, label, color in nodes:
        draw.line((center[0], center[1], x, y), fill=rgb("#4c7895"), width=4)
        draw.rounded_rectangle((x - 82, y - 47, x + 82, y + 47), radius=18, fill=rgb("#102f4f"), outline=rgb(color), width=4)
        draw.rectangle((x - 38, y - 18, x + 38, y + 18), outline=rgb(color), width=4)
        draw.line((x - 24, y - 6, x + 24, y - 6), fill=rgb(color), width=3)
        draw.text((x - 34, y + 58), label, fill=rgb("#f2f8ff"))
    draw.text((321, 445), "3 replicas  |  2 matching copies required", fill=rgb("#b8cfe5"))
    return image_bytes(canvas)


def make_evidence_visual() -> io.BytesIO:
    """Create a file/hash illustration for the evidence-fixtures slide."""
    canvas = Image.new("RGB", (900, 520), rgb("#081c35"))
    draw = ImageDraw.Draw(canvas)
    files = [(70, 92, "TXT", "camera-note"), (300, 92, "JSON", "mobile-extract"), (530, 92, "CSV", "firewall-log"), (185, 290, "LOG", "browser-session"), (415, 290, "HTML", "case-report")]
    colors = ["#63cdff", "#53dcb5", "#f6b859", "#f77e89", "#b68bff"]
    for (x, y, kind, name), color in zip(files, colors):
        draw.rounded_rectangle((x, y, x + 190, y + 120), radius=16, fill=rgb("#102f4f"), outline=rgb(color), width=4)
        draw.rectangle((x + 18, y + 20, x + 64, y + 82), outline=rgb(color), width=4)
        draw.line((x + 29, y + 40, x + 53, y + 40), fill=rgb(color), width=3)
        draw.line((x + 29, y + 55, x + 53, y + 55), fill=rgb(color), width=3)
        draw.text((x + 78, y + 22), kind, fill=rgb(color))
        draw.text((x + 78, y + 61), name, fill=rgb("#f2f8ff"))
    draw.line((110, 470, 790, 470), fill=rgb("#3d789c"), width=4)
    draw.ellipse((432, 440, 468, 476), fill=rgb("#53dcb5"))
    draw.text((175, 488), "SHA-256 fingerprint  →  distributed replicas  →  custody report", fill=rgb("#b8cfe5"))
    return image_bytes(canvas)


def add_background(slide: object, bg: io.BytesIO) -> None:
    bg.seek(0)
    slide.shapes.add_picture(bg, 0, 0, width=W, height=H)


def fill_shape(shape: object, color: RGBColor, line: RGBColor | None = None) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = line or color


def add_box(slide: object, x: float, y: float, w: float, h: float, color: RGBColor = PANEL, line: RGBColor = PANEL_2, radius: bool = True) -> object:
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    fill_shape(shape, color, line)
    if radius:
        shape.adjustments[0] = 0.12
    return shape


def add_text(slide: object, text: str, x: float, y: float, w: float, h: float, size: float = 18, color: RGBColor = WHITE, bold: bool = False, font: str = "Aptos", align: PP_ALIGN = PP_ALIGN.LEFT, margin: float = 0.04, wrap: bool = True) -> object:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = wrap
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = PRESENTATION_FONT
    # Slightly enlarge small presentation text so it remains readable when projected.
    readable_size = size * 1.20 if size < 20 else size
    run.font.size = Pt(readable_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_label(slide: object, text: str, x: float, y: float, w: float, color: RGBColor = CYAN) -> None:
    add_text(slide, text.upper(), x, y, w, 0.26, 9, color, True, "Aptos Display")


def add_title(slide: object, title: str, subtitle: str, number: str) -> None:
    add_label(slide, f"CUSTODIA  /  {number}", 0.62, 0.42, 4.2)
    add_text(slide, title, 0.62, 0.78, 8.7, 0.62, 27, WHITE, True)
    add_text(slide, subtitle, 0.64, 1.42, 10.8, 0.38, 12, MUTED)
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.64), Inches(1.92), Inches(12.68), Inches(1.92))
    line.line.color.rgb = RGBColor(65, 115, 158)
    line.line.width = Pt(1.2)


def add_footer(slide: object, text: str = "GROUP-I  |  CUSTODIA EVIDENCE NETWORK", number: str = "01") -> None:
    add_text(slide, text, 0.64, 7.16, 7.5, 0.2, 8, MUTED, True, "Aptos Display")
    add_text(slide, number, 12.12, 7.16, 0.55, 0.2, 8, CYAN, True, "Aptos Display", PP_ALIGN.RIGHT)


def add_bullet_list(slide: object, items: list[str], x: float, y: float, w: float, line_height: float = 0.55, size: float = 15, color: RGBColor = WHITE) -> None:
    for index, item in enumerate(items):
        yy = y + index * line_height
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(yy + 0.14), Inches(0.11), Inches(0.11))
        fill_shape(dot, CYAN, CYAN)
        add_text(slide, item, x + 0.25, yy, w - 0.25, line_height, size, color)


def add_step(slide: object, number: str, title: str, body: str, x: float, y: float, w: float, accent: RGBColor = CYAN) -> None:
    card = add_box(slide, x, y, w, 1.42, PANEL, RGBColor(48, 100, 145))
    badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.18), Inches(y + 0.22), Inches(0.48), Inches(0.48))
    fill_shape(badge, accent, accent)
    add_text(slide, number, x + 0.18, y + 0.22, 0.48, 0.48, 12, NAVY, True, "Aptos Display", PP_ALIGN.CENTER)
    add_text(slide, title, x + 0.82, y + 0.19, w - 1.02, 0.34, 14, WHITE, True)
    add_text(slide, body, x + 0.82, y + 0.58, w - 1.02, 0.62, 10.5, MUTED)


def build() -> None:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]
    bg = make_background()

    # Slide 1: title
    slide = prs.slides.add_slide(blank)
    add_background(slide, bg)
    add_label(slide, "DIGITAL FORENSICS  /  PROJECT PRESENTATION", 0.78, 0.78, 6.2, TEAL)
    add_text(slide, "CUSTODIA", 0.72, 1.45, 7.4, 0.92, 47, WHITE, True, "Aptos Display")
    add_text(slide, "Distributed Digital Evidence\nand Chain of Custody", 0.76, 2.48, 7.1, 1.14, 25, CYAN, True)
    add_text(slide, "A local forensic evidence network for integrity, redundancy, recovery, and accountability.", 0.8, 3.86, 6.6, 0.6, 14, MUTED)
    accent = add_box(slide, 0.8, 5.18, 3.3, 0.68, PANEL_2, RGBColor(87, 180, 225))
    add_text(slide, "PRESENTED BY  GROUP-I", 1.04, 5.18, 2.82, 0.68, 16, WHITE, True, "Aptos Display", PP_ALIGN.CENTER)
    add_text(slide, "Forensics theme  |  SHA-256  |  Quorum storage", 0.82, 6.46, 6.5, 0.28, 10, MUTED, True)

    # Slide 2: members
    slide = prs.slides.add_slide(blank)
    add_background(slide, bg)
    add_title(slide, "Group Members", "Six-member project team with individual names and identification codes.", "02")
    positions = [(0.72, 2.28), (4.54, 2.28), (8.36, 2.28), (0.72, 4.56), (4.54, 4.56), (8.36, 4.56)]
    members = [
        ("REAL NAME 01", "TNT-1963"),
        ("REAL NAME 02", "TNT-1964"),
        ("REAL NAME 03", "TNT-1965"),
        ("REAL NAME 04", "TNT-1966"),
        ("REAL NAME 05", "TNT-1967"),
        ("REAL NAME 06", "TNT-1968"),
    ]
    for index, (x, y) in enumerate(positions):
        member_name, member_id = members[index]
        card = add_box(slide, x, y, 3.32, 1.84, PANEL, RGBColor(48, 100, 145))
        avatar = avatar_bytes(index, index >= 3)
        avatar.seek(0)
        slide.shapes.add_picture(avatar, Inches(x + 0.14), Inches(y + 0.25), width=Inches(1.12), height=Inches(1.12))
        add_text(slide, member_name, x + 1.34, y + 0.36, 1.78, 0.42, 14, CYAN, True, wrap=False)
        add_text(slide, f"ID: {member_id}", x + 1.34, y + 0.96, 1.78, 0.28, 12.5, TEAL, True, wrap=False)
    add_footer(slide, number="02")

    # Slide 3: agenda
    slide = prs.slides.add_slide(blank)
    add_background(slide, bg)
    add_title(slide, "Agenda", "A simple route from the problem to a working forensic evidence demonstration.", "03")
    add_box(slide, 0.82, 2.35, 6.2, 3.95, PANEL, RGBColor(48, 100, 145))
    add_label(slide, "TODAY'S ROUTE", 1.12, 2.7, 2.2, TEAL)
    agenda = [
        "Project purpose and learning objectives",
        "Forensic evidence tools and system architecture",
        "Upload, hash, replicate, verify, and repair",
        "Testing with fictional evidence files",
        "Benefits, limitations, and future direction",
        "Conclusion and questions",
    ]
    add_bullet_list(slide, agenda, 1.12, 3.16, 5.35, 0.49, 14, WHITE)
    tools_visual = make_tools_visual()
    tools_visual.seek(0)
    slide.shapes.add_picture(tools_visual, Inches(7.45), Inches(2.45), width=Inches(5.1), height=Inches(2.95))
    add_box(slide, 7.45, 5.58, 5.1, 0.62, RGBColor(11, 36, 64), RGBColor(80, 158, 193))
    add_text(slide, "Understand the workflow, then run the demo", 7.68, 5.58, 4.64, 0.62, 13, WHITE, True, "Aptos Display", PP_ALIGN.CENTER)
    add_footer(slide, number="03")

    # Slide 4: overview and objectives
    slide = prs.slides.add_slide(blank)
    add_background(slide, bg)
    add_title(slide, "Project overview and objectives", "Custodia is a learning prototype for protecting digital evidence from intake to report.", "04")
    cards = [
        (0.82, "CAPTURE", "Register evidence with a case ID, description, and tags.", CYAN),
        (3.98, "PROTECT", "Hash the file and store redundant copies across nodes.", TEAL),
        (7.14, "PROVE", "Verify integrity and show who performed each action.", AMBER),
        (10.3, "RECOVER", "Detect corruption, repair replicas, and preserve history.", RED),
    ]
    for card_index, (x, label, body, accent) in enumerate(cards, start=1):
        add_box(slide, x, 2.55, 2.25, 2.55, PANEL, RGBColor(48, 100, 145))
        add_box(slide, x + 0.2, 2.82, 0.52, 0.52, accent, accent)
        add_text(slide, str(card_index), x + 0.2, 2.82, 0.52, 0.52, 14, NAVY, True, "Aptos Display", PP_ALIGN.CENTER)
        add_label(slide, label, x + 0.2, 3.56, 1.65, accent)
        add_text(slide, body, x + 0.2, 3.98, 1.82, 0.76, 11.5, WHITE)
    add_box(slide, 1.65, 5.55, 10.0, 0.65, RGBColor(11, 36, 64), RGBColor(80, 158, 193))
    add_text(slide, "Main lesson: evidence is not only a file—it is a file plus proof, redundancy, and history.", 1.9, 5.55, 9.5, 0.65, 14, WHITE, True, "Aptos Display", PP_ALIGN.CENTER)
    add_footer(slide, number="04")

    # Slide 5: problem
    slide = prs.slides.add_slide(blank)
    add_background(slide, bg)
    add_title(slide, "Why digital evidence needs protection", "Evidence is valuable only when its origin, content, and history can be trusted.", "05")
    add_box(slide, 0.72, 2.35, 5.65, 3.9, PANEL, RGBColor(48, 100, 145))
    add_label(slide, "THE RISK", 1.02, 2.7, 2.0, RED)
    add_text(slide, "A single file copy is not enough", 1.0, 3.02, 4.7, 0.5, 22, WHITE, True)
    add_bullet_list(slide, ["Files can be changed without obvious signs.", "A failed drive can make evidence unavailable.", "Manual history is easy to miss or dispute.", "Different versions can be confused."], 1.02, 3.75, 4.7, 0.59, 14, MUTED)
    add_box(slide, 6.72, 2.35, 5.85, 3.9, PANEL_2, RGBColor(80, 158, 193))
    add_label(slide, "THE RESPONSE", 7.02, 2.7, 2.4, TEAL)
    add_text(slide, "Make every byte accountable", 7.0, 3.02, 4.9, 0.5, 22, WHITE, True)
    add_bullet_list(slide, ["Create a cryptographic fingerprint.", "Keep three replicas across four nodes.", "Verify reads with a majority quorum.", "Record every action in a linked ledger."], 7.02, 3.75, 4.95, 0.59, 14, WHITE)
    add_footer(slide, number="05")

    # Slide 6: forensic tools
    slide = prs.slides.add_slide(blank)
    add_background(slide, bg)
    add_title(slide, "Forensic tools inside the project", "Each tool answers a different trust question during evidence handling.", "06")
    visual = make_tools_visual()
    visual.seek(0)
    slide.shapes.add_picture(visual, Inches(0.82), Inches(2.35), width=Inches(6.3), height=Inches(3.65))
    add_box(slide, 7.45, 2.35, 5.1, 3.65, PANEL, RGBColor(48, 100, 145))
    add_label(slide, "WHAT EACH TOOL PROVES", 7.75, 2.72, 3.2, TEAL)
    add_bullet_list(slide, ["File intake: what was submitted?", "SHA-256: did the bytes change?", "Node replicas: is the file still available?", "Quorum read: do copies agree?", "Audit ledger: who did what and when?"], 7.75, 3.18, 4.35, 0.52, 13, WHITE)
    add_footer(slide, number="06")

    # Slide 7: architecture
    slide = prs.slides.add_slide(blank)
    add_background(slide, bg)
    add_title(slide, "Project architecture", "The browser dashboard calls FastAPI, which manages metadata, hashes, replicas, and audit history.", "07")
    add_box(slide, 0.82, 2.38, 2.45, 2.25, PANEL_2, RGBColor(80, 158, 193))
    add_label(slide, "FRONTEND", 1.08, 2.7, 1.4, CYAN)
    add_text(slide, "React + TypeScript", 1.05, 3.12, 1.95, 0.5, 18, WHITE, True)
    add_text(slide, "Dashboard\nUpload, verify, repair, report", 1.08, 3.78, 1.85, 0.62, 11, MUTED)
    add_box(slide, 5.0, 2.38, 3.0, 2.25, PANEL, RGBColor(48, 100, 145))
    add_label(slide, "BACKEND", 5.3, 2.7, 1.4, TEAL)
    add_text(slide, "Python + FastAPI", 5.28, 3.12, 2.35, 0.5, 18, WHITE, True)
    add_text(slide, "Hashing\nQuorum reads\nVersioning + permissions", 5.3, 3.73, 2.3, 0.78, 11, MUTED)
    add_box(slide, 9.75, 2.38, 2.75, 2.25, PANEL, RGBColor(48, 100, 145))
    add_label(slide, "LOCAL DATA", 10.05, 2.7, 1.8, AMBER)
    add_text(slide, "metadata.json", 10.02, 3.12, 2.2, 0.4, 16, WHITE, True)
    add_text(slide, "Four node folders\nObject chunks\nQuarantine copies", 10.05, 3.73, 2.1, 0.78, 11, MUTED)
    for x1, x2 in [(3.35, 4.9), (8.05, 9.65)]:
        arrow = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(3.5), Inches(x2), Inches(3.5))
        arrow.line.color.rgb = CYAN
        arrow.line.width = Pt(2.2)
        arrow.line.end_arrowhead = True
    add_box(slide, 2.32, 5.25, 8.7, 0.7, RGBColor(11, 36, 64), RGBColor(48, 100, 145))
    add_text(slide, "User action  →  API route  →  evidence service  →  verified result + audit event", 2.55, 5.25, 8.22, 0.7, 15, WHITE, True, "Aptos Display", PP_ALIGN.CENTER)
    add_footer(slide, number="07")

    # Slide 8: workflow
    slide = prs.slides.add_slide(blank)
    add_background(slide, bg)
    add_title(slide, "End-to-end evidence workflow", "One upload moves through intake, distribution, verification, recovery, and reporting.", "08")
    steps = [
        ("01", "Register", "Choose a file, case ID, description, and tags.", CYAN),
        ("02", "Fingerprint", "SHA-256 hashes the file and every chunk.", TEAL),
        ("03", "Distribute", "Three replicas are written to the node ring.", AMBER),
        ("04", "Verify", "Two matching replicas are required for a read.", CYAN),
        ("05", "Recover", "Bad copies are quarantined and rebuilt.", RED),
        ("06", "Report", "Versions, hashes, events, and certification export.", TEAL),
    ]
    for index, (num, title, body, accent) in enumerate(steps):
        x = 0.82 + (index % 3) * 4.15
        y = 2.4 + (index // 3) * 1.95
        add_step(slide, num, title, body, x, y, 3.62, accent)
    add_footer(slide, number="08")

    # Slide 9: integrity model
    slide = prs.slides.add_slide(blank)
    add_background(slide, bg)
    add_title(slide, "Integrity, replication, and recovery", "Custodia validates both the individual chunks and the reconstructed whole file.", "09")
    add_box(slide, 0.78, 2.35, 7.05, 3.95, PANEL, RGBColor(48, 100, 145))
    add_label(slide, "STORAGE MODEL", 1.08, 2.7, 2.2, CYAN)
    add_text(slide, "One file becomes trusted building blocks", 1.06, 3.06, 5.7, 0.48, 20, WHITE, True)
    nodes = [("FILE", 1.1, 4.2, CYAN), ("CHUNKS", 3.05, 4.2, TEAL), ("3 REPLICAS", 5.0, 4.2, AMBER)]
    for label, x, y, accent in nodes:
        circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(1.15), Inches(1.15))
        fill_shape(circle, accent, accent)
        add_text(slide, label, x, y + 0.35, 1.15, 0.35, 10, NAVY, True, "Aptos Display", PP_ALIGN.CENTER)
    for x1, x2 in [(2.25, 2.95), (4.2, 4.9)]:
        arrow = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(4.78), Inches(x2), Inches(4.78))
        arrow.line.color.rgb = MUTED
        arrow.line.width = Pt(2)
        arrow.line.end_arrowhead = True
    add_text(slide, "File hash", 1.16, 5.55, 1.05, 0.25, 10, MUTED, True, "Aptos Display", PP_ALIGN.CENTER)
    add_text(slide, "Merkle root", 3.06, 5.55, 1.15, 0.25, 10, MUTED, True, "Aptos Display", PP_ALIGN.CENTER)
    add_text(slide, "2 of 3 quorum", 5.0, 5.55, 1.18, 0.25, 10, MUTED, True, "Aptos Display", PP_ALIGN.CENTER)
    add_box(slide, 8.2, 2.35, 4.35, 4.2, PANEL_2, RGBColor(80, 158, 193))
    add_label(slide, "WHEN A COPY CHANGES", 8.5, 2.7, 3.1, RED)
    add_text(slide, "Detect → quarantine → repair", 8.48, 3.06, 3.6, 0.48, 20, WHITE, True)
    add_bullet_list(slide, ["Compare bytes with the expected hash.", "Use a healthy peer for repair.", "Quarantine the bad copy.", "Write correct bytes atomically.", "Record repair in the ledger."], 8.5, 3.78, 3.5, 0.5, 10.5, WHITE)
    add_footer(slide, number="09")

    # Slide 10: evidence files and case IDs
    slide = prs.slides.add_slide(blank)
    add_background(slide, bg)
    add_title(slide, "Evidence files and case IDs", "Different file types can belong to the same case without becoming the same evidence item.", "10")
    evidence_visual = make_evidence_visual()
    evidence_visual.seek(0)
    slide.shapes.add_picture(evidence_visual, Inches(0.72), Inches(2.35), width=Inches(5.92), height=Inches(3.67))
    add_box(slide, 6.92, 2.35, 5.63, 4.2, PANEL, RGBColor(48, 100, 145))
    add_label(slide, "BEGINNER RULE", 7.22, 2.72, 2.4, AMBER)
    add_text(slide, "One case can contain many files", 7.2, 3.12, 4.95, 0.5, 18, WHITE, True)
    add_bullet_list(slide, ["Case ID: CASE-TEST-001", "Evidence ID: created per upload", "Same case ID groups related files", "Filenames keep versions distinct", "Example: camera note + mobile extract"], 7.22, 3.78, 4.95, 0.48, 11.5, WHITE)
    add_footer(slide, number="10")

    # Slide 11: testing
    slide = prs.slides.add_slide(blank)
    add_background(slide, bg)
    add_title(slide, "Testing the project", "The supplied fictional evidence files let the team demonstrate every current feature safely.", "11")
    add_box(slide, 0.78, 2.35, 5.25, 3.95, PANEL, RGBColor(48, 100, 145))
    add_label(slide, "SAMPLE FILES", 1.08, 2.7, 2.0, CYAN)
    add_text(slide, "21 fictional fixtures", 1.06, 3.06, 3.6, 0.45, 22, WHITE, True)
    add_bullet_list(slide, ["TXT: camera and interview notes", "JSON: mobile and browser extracts", "CSV: firewall, GPS, USB, timelines", "LOG, Markdown, and HTML reports", "Use CASE-TEST-001 to group related files"], 1.08, 3.78, 4.3, 0.47, 13, MUTED)
    add_box(slide, 6.45, 2.35, 6.1, 3.95, PANEL_2, RGBColor(80, 158, 193))
    add_label(slide, "LIVE TEST SEQUENCE", 6.75, 2.7, 2.4, TEAL)
    tests = [("1", "Upload", "Hash & distribute"), ("2", "Verify", "Healthy result"), ("3", "Version", "v1 stays immutable"), ("4", "Corrupt", "Status becomes attention"), ("5", "Repair", "Replica returns healthy"), ("6", "Report", "JSON custody export")]
    for index, (num, title, body) in enumerate(tests):
        x = 6.75 + (index % 2) * 2.8
        y = 3.2 + (index // 2) * 0.82
        add_text(slide, num, x, y, 0.28, 0.3, 13, CYAN, True, "Aptos Display", PP_ALIGN.CENTER)
        add_text(slide, title, x + 0.38, y, 1.25, 0.28, 12, WHITE, True)
        add_text(slide, body, x + 0.38, y + 0.29, 2.05, 0.22, 9.5, MUTED)
    add_footer(slide, number="11")

    # Slide 12: roles and use cases
    slide = prs.slides.add_slide(blank)
    add_background(slide, bg)
    add_title(slide, "Who uses Custodia?", "The prototype makes responsibilities visible for investigators, examiners, and administrators.", "12")
    use_cases = [
        (0.82, "INVESTIGATOR", "Register new evidence, add notes, create versions, and follow case activity.", CYAN),
        (4.25, "EXAMINER", "Download files, check hashes, compare versions, and produce a custody report.", TEAL),
        (7.68, "ADMINISTRATOR", "Repair replicas, simulate faults, review audit events, and manage system health.", AMBER),
    ]
    for x, role, body, accent in use_cases:
        add_box(slide, x, 2.55, 3.05, 2.95, PANEL, RGBColor(48, 100, 145))
        badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.26), Inches(2.86), Inches(0.66), Inches(0.66))
        fill_shape(badge, accent, accent)
        add_text(slide, role[0], x + 0.26, 2.86, 0.66, 0.66, 17, NAVY, True, "Aptos Display", PP_ALIGN.CENTER)
        add_label(slide, role, x + 1.08, 2.98, 1.8, accent)
        add_text(slide, body, x + 0.28, 3.9, 2.46, 0.9, 12.5, WHITE)
        add_text(slide, "Permission depends on role", x + 0.28, 5.03, 2.45, 0.25, 10, MUTED, True)
    add_box(slide, 2.15, 5.92, 8.95, 0.54, RGBColor(11, 36, 64), RGBColor(80, 158, 193))
    add_text(slide, "Use fictional files for demonstrations; never upload real sensitive evidence to this prototype.", 2.38, 5.92, 8.5, 0.54, 12, AMBER, True, "Aptos Display", PP_ALIGN.CENTER)
    add_footer(slide, number="12")

    # Slide 13: benefits and limits
    slide = prs.slides.add_slide(blank)
    add_background(slide, bg)
    add_title(slide, "Benefits today, improvements tomorrow", "Custodia is a strong learning prototype with a clear path toward production infrastructure.", "13")
    add_box(slide, 0.78, 2.35, 5.9, 3.95, PANEL, RGBColor(48, 100, 145))
    add_label(slide, "CURRENT BENEFITS", 1.08, 2.7, 2.5, TEAL)
    add_bullet_list(slide, ["Tamper detection with SHA-256", "Redundant storage and quorum reads", "Automatic corruption recovery", "Immutable version history", "Role-based actions and audit reports"], 1.08, 3.18, 5.0, 0.55, 14, WHITE)
    add_box(slide, 7.0, 2.35, 5.55, 3.95, PANEL_2, RGBColor(80, 158, 193))
    add_label(slide, "NEXT ADVANCED STEPS", 7.3, 2.7, 2.8, AMBER)
    add_bullet_list(slide, ["PostgreSQL + database migrations", "MinIO/S3 storage on separate servers", "Real login, MFA, and JWT/OIDC", "AES-256 encryption with KMS", "Scheduled workers, alerts, OCR, and analytics"], 7.3, 3.18, 4.6, 0.55, 14, WHITE)
    add_text(slide, "Important: the current nodes are local folders and users are demo identities.", 1.08, 6.58, 10.9, 0.32, 11, AMBER, True)
    add_footer(slide, number="13")

    # Slide 14: roadmap
    slide = prs.slides.add_slide(blank)
    add_background(slide, bg)
    add_title(slide, "Roadmap to an advanced platform", "The current prototype demonstrates the workflow; these upgrades prepare it for real organizations.", "14")
    roadmap = [
        ("NOW", "Prototype", "metadata.json, local node folders, demo identities", CYAN),
        ("NEXT", "Reliable service", "PostgreSQL, migrations, MinIO/S3, background jobs", TEAL),
        ("LATER", "Production security", "AES-256, KMS, TLS, MFA, signed transfers", AMBER),
        ("GOAL", "Operational platform", "OCR, malware scanning, alerts, reports, retention", RED),
    ]
    for index, (phase, title, body, accent) in enumerate(roadmap):
        x = 0.82 + index * 3.1
        add_box(slide, x, 2.7, 2.7, 2.45, PANEL, RGBColor(48, 100, 145))
        add_label(slide, phase, x + 0.25, 3.02, 1.0, accent)
        add_text(slide, title, x + 0.25, 3.45, 2.1, 0.4, 17, WHITE, True)
        add_text(slide, body, x + 0.25, 4.02, 2.1, 0.72, 11, MUTED)
        if index < len(roadmap) - 1:
            arrow = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + 2.7), Inches(3.92), Inches(x + 3.02), Inches(3.92))
            arrow.line.color.rgb = CYAN
            arrow.line.width = Pt(2)
            arrow.line.end_arrowhead = True
    add_box(slide, 2.05, 5.72, 9.2, 0.58, RGBColor(11, 36, 64), RGBColor(80, 158, 193))
    add_text(slide, "Build one trust layer at a time: data → storage → security → automation", 2.28, 5.72, 8.75, 0.58, 13, WHITE, True, "Aptos Display", PP_ALIGN.CENTER)
    add_footer(slide, number="14")

    # Slide 15: closing
    slide = prs.slides.add_slide(blank)
    add_background(slide, bg)
    add_label(slide, "GROUP-I  /  CONCLUSION", 0.82, 1.18, 3.0, TEAL)
    add_text(slide, "Every byte accounted for.\nEvery action attributable.", 0.78, 2.0, 8.6, 1.25, 32, WHITE, True)
    add_text(slide, "Custodia demonstrates how hashing, replication, recovery, and audit history work together to protect digital evidence.", 0.82, 3.62, 7.8, 0.66, 16, MUTED)
    add_text(slide, "Questions and discussion", 0.82, 4.55, 4.8, 0.5, 19, CYAN, True)
    add_box(slide, 0.82, 5.1, 4.3, 0.72, PANEL, RGBColor(80, 158, 193))
    add_text(slide, "CUSTODIA  |  DIGITAL FORENSICS", 1.08, 5.1, 3.78, 0.72, 14, WHITE, True, "Aptos Display", PP_ALIGN.CENTER)
    add_footer(slide, number="15")

    prs.core_properties.title = "Custodia - Distributed Digital Evidence"
    prs.core_properties.subject = "Group-I forensic evidence project presentation"
    prs.core_properties.author = "Group-I"
    prs.core_properties.keywords = "digital forensics, chain of custody, SHA-256, quorum storage"
    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    build()
